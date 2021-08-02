from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN


#データの定義
day = "8/3火"
menu_array = []
menu_array.append("DXトンカツ弁当")
menu_array.append("DXメンチカツ弁当")
menu_array.append("ハンバーグカレー")
menu_array.append("鶏チャーハン")

#プレゼンテーションを開く
prs = Presentation()

prs.slide_width = Cm(27.52)
prs.slide_height = Cm(19.05)

#1ページ目（「タイトル スライド」のレイアウトを指定）
slide_layout_1 = prs.slide_layouts[1]
slide_1 = prs.slides.add_slide(slide_layout_1)

#プレースホルダーの指定
s1_title = slide_1.placeholders[0]
s1_contents = slide_1.placeholders[1]

#textをそれぞれのプレースフォルダに設定
s1_title.text = "今日のお弁当(" + day + ")"

menu_text = ""
for menu in menu_array:
    menu_text = menu_text + menu + "\n"
s1_contents.text = menu_text

slide_1.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

prs.save("todays_menu.pptx")