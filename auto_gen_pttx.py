from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt
#フォントの色を変えるためにRGBcolorを使う
from pptx.dml.color import RGBColor
#テキストの配置を調整する時にMSO_ANCHORとPP_ALIGNを使う
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

def auto_gen(array):

    #データの定義
    day = "8/3火"
    menu_array = array
    # menu_array.append("DXトンカツ弁当")
    # menu_array.append("DXメンチカツ弁当")
    # menu_array.append("ハンバーグカレー")
    # menu_array.append("鶏チャーハン")
    #テストコメント

    #プレゼンテーションを開く
    prs = Presentation()

    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)

    #1ページ目（「タイトル スライド」のレイアウトを指定）
    slide_layout_6 = prs.slide_layouts[6]
    slide_1 = prs.slides.add_slide(slide_layout_6)

    # テキストボックスを追加
    shapes = slide_1.shapes
    shape = shapes.add_textbox(Cm(1), Cm(1), Cm(26), Cm(16))

    #shapeオブジェクトでtextを挿入
    text_frame = shape.text_frame
    #パラグラフの追加
    pg = text_frame.paragraphs[0]
    #フォントサイズを大きくする
    pg.font.italic = True
    pg.font.bold = True
    pg.font.size = Pt(60)
    pg.text = "今日のお弁当(" + day + ")"


    #以下、フォントをMeiryo UIにしてboldにする
    for index, menu in enumerate(menu_array):
        pg = text_frame.add_paragraph()
        pg.font.bold = True
        pg.font.size = Pt(38)
        # pg.font.name = "font.name = Meiryo UI"
        pg.text = str(index+1) + ". " + menu

    shape = shapes.add_textbox(Cm(1), Cm(16.5), Cm(18), Cm(2))
    text_frame = shape.text_frame
    pg = text_frame.paragraphs[0]
    pg.font.size = Pt(24)
    pg.text = "※カッコ書きが無いお弁当は、550円"

    shape = shapes.add_textbox(Cm(18.5), Cm(16.5), Cm(7), Cm(2))
    text_frame = shape.text_frame
    pg = text_frame.paragraphs[0]
    pg.font.size = Pt(32)
    run = pg.add_run()
    run.text = "笑"
    run = pg.add_run()
    run.text = "楽"
    run.font.color.rgb = RGBColor(255, 0, 0)
    run = pg.add_run()
    run.text = "のお弁当"

    #以下、フォントをイタリックにするにする
    # pg = text_frame.add_paragraph()
    # pg.font.italic = True
    # pg.text = "font.italic = True"

    #以下、アンダーラインを引く
    # pg = text_frame.add_paragraph()
    # pg.font.underline = True
    # pg.text = "font.underline = True"


    #プレースホルダーの指定
    # s1_title = slide_1.placeholders[0]
    # s1_contents = slide_1.placeholders[1]

    #textをそれぞれのプレースフォルダに設定
    # s1_title.text = "今日のお弁当(" + day + ")"

    # menu_text = ""
    # for index, menu in enumerate(menu_array):
    #     menu_text =  menu_text + str(index+1) + ". " + menu + "\n"
    # s1_contents.text = menu_text


    # slide_1.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    prs.save("todays_menu.pptx")