from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt
#フォントの色を変えるためにRGBcolorを使う
from pptx.dml.color import RGBColor
#テキストの配置を調整する時にMSO_ANCHORとPP_ALIGNを使う
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import datetime
import locale

import my_util

def auto_gen(date_str, menu_array):
    locale.setlocale(locale.LC_TIME, 'ja_JP.UTF-8')
    
    file_path = '../output/todays_menu.pptx'
    datetime_date = my_util.get_datetime(date_str)

    #プレゼンテーションを開く
    prs = Presentation()

    prs.slide_width = Cm(27.52)
    prs.slide_height = Cm(19.05)

    #1ページ目（「タイトル スライド」のレイアウトを指定）
    slide_layout_6 = prs.slide_layouts[6]
    slide_1 = prs.slides.add_slide(slide_layout_6)

    shapes = slide_1.shapes

    # テキストボックスを追加
    shape = shapes.add_textbox(Cm(0.3), Cm(0.3), Cm(26), Cm(3))

    #shapeオブジェクトでtextを挿入
    text_frame = shape.text_frame
    #パラグラフの追加
    pg = text_frame.paragraphs[0]
    #フォントサイズを大きくする
    pg.font.italic = True
    pg.font.bold = True
    pg.font.size = Pt(60)
    pg.text = "今日のお弁当(" + date_str + datetime_date.strftime('%a') + ")"

    # テキストボックスを追加
    shape = shapes.add_textbox(Cm(1), Cm(3.5), Cm(26), Cm(13))
    text_frame = shape.text_frame

    #メニューの行数によって文字の大きさを設定する
    if len(menu_array) <= 6 :
        font_size = Pt(48)
    elif len(menu_array) <= 8 :
        font_size = Pt(38)
    else :
        font_size = Pt(30)

    #以下、フォントをMeiryo UIにしてboldにする
    for index, menu in enumerate(menu_array):
        if index == 0:
            pg = text_frame.paragraphs[0]
        else :
            pg = text_frame.add_paragraph()
        pg.font.bold = True
        
        pg.font.size = font_size
        # pg.font.name = "font.name = Meiryo UI"
        pg.text = str(index+1) + ". " + menu


    #注意書きと店名
    shape = shapes.add_textbox(Cm(1), Cm(16.5), Cm(18), Cm(2))
    text_frame = shape.text_frame
    pg = text_frame.paragraphs[0]
    pg.font.size = Pt(24)
    pg.text = "※カッコ書きが無いお弁当は、550円"

    shape = shapes.add_textbox(Cm(18.5), Cm(16.5), Cm(7), Cm(2))
    text_frame = shape.text_frame
    pg = text_frame.paragraphs[0]
    pg.font.size = Pt(32)
    run = pg.add_run()
    run.text = "笑"
    run = pg.add_run()
    run.text = "楽"
    run.font.color.rgb = RGBColor(255, 0, 0)
    run = pg.add_run()
    run.text = "のお弁当"

    prs.save(file_path)